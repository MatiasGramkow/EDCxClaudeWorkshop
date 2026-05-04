"""Generate session 1 PowerPoint deck for the EDC × Claude Code workshop.

Designet matcher EDC's officielle workshop-template:
  - Hvid baggrund med tynd navy-stribe top og bund
  - "**Bold** | Lighter"-titler i navy
  - Lyse blågrå indholds-bokse
  - Sektions-dividers med fuld navy + 'Punkt N' i gult
  - EDC-logo nederst til højre på alle indholds-slides
  - Sidetal nederst til venstre (format: "N |")

Output: session1.pptx (16:9). Åbn i Keynote og brug Presenter Display
til at vise slides på projektor og noter på din Mac.

Run: python3 build_session1.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH = os.path.join(HERE, 'edc_logo.png')

# ---------- EDC palette --------------------------------------------------
EDC_NAVY = RGBColor(0x12, 0x2A, 0x52)        # primary navy (matches PDF)
EDC_NAVY_DARK = RGBColor(0x0A, 0x1F, 0x3D)
EDC_BLUE_BODY = RGBColor(0x1F, 0x3A, 0x6B)   # body text navy
EDC_YELLOW = RGBColor(0xF5, 0xB8, 0x1C)      # accent yellow ("Punkt N")
EDC_LIGHT_BG = RGBColor(0xEE, 0xF1, 0xF6)    # content box background
EDC_GREY = RGBColor(0x9C, 0xA3, 0xAF)
INK = RGBColor(0x12, 0x2A, 0x52)
MUTED = RGBColor(0x55, 0x6B, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_BG = RGBColor(0xD7, 0xEC, 0xDD)
GREEN_INK = RGBColor(0x14, 0x53, 0x2D)
RED_BG = RGBColor(0xFA, 0xDC, 0xDC)
RED_INK = RGBColor(0x8B, 0x1F, 0x1F)
CODE_BG = RGBColor(0x0F, 0x1F, 0x3D)
CODE_INK = RGBColor(0xE2, 0xE8, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------- Low-level helpers ---------------------------------------------
def fill_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(slide, left, top, width, height, text, *,
             font='Helvetica', size=24, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tx


def _parse_bold(text):
    """Split 'foo **bar** baz' into [('foo ', False), ('bar', True), (' baz', False)]"""
    out = []
    i = 0
    while i < len(text):
        if text[i:i+2] == '**':
            end = text.find('**', i + 2)
            if end == -1:
                out.append((text[i:], False))
                break
            out.append((text[i+2:end], True))
            i = end + 2
        else:
            nxt = text.find('**', i)
            if nxt == -1:
                out.append((text[i:], False))
                break
            out.append((text[i:nxt], False))
            i = nxt
    return [seg for seg in out if seg[0]]


def add_rich_text(slide, left, top, width, height, text, *,
                  font='Helvetica', size=22, color=INK,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                  line_spacing=1.3):
    """Single-paragraph text with **bold** support."""
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for seg, is_bold in _parse_bold(text):
        run = p.add_run()
        run.text = seg
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = is_bold
        run.font.color.rgb = color
    return tx


def add_bullets(slide, left, top, width, height, bullets, *,
                font='Helvetica', size=20, color=INK, line_spacing=1.4,
                bullet_glyph='■', bullet_color=None):
    if bullet_color is None:
        bullet_color = EDC_NAVY
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(10)
        # bullet glyph
        br = p.add_run()
        br.text = f'{bullet_glyph}   '
        br.font.name = font
        br.font.size = Pt(size)
        br.font.bold = True
        br.font.color.rgb = bullet_color
        for seg, is_bold in _parse_bold(b):
            run = p.add_run()
            run.text = seg
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = is_bold
            run.font.color.rgb = color
    return tx


def add_code_block(slide, left, top, width, height, code, *, size=15):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.04
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_right = Inches(0.4)
    tf.margin_top = Inches(0.3)
    tf.margin_bottom = Inches(0.3)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = code.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        run = p.add_run()
        run.text = line if line else ' '
        run.font.name = 'Menlo'
        run.font.size = Pt(size)
        run.font.color.rgb = CODE_INK
    return box


# ---------- EDC chrome ----------------------------------------------------
def add_chrome(slide, page_num):
    """Top + bottom navy stripes, EDC logo, page number."""
    # top stripe
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    top.fill.solid()
    top.fill.fore_color.rgb = EDC_NAVY
    top.line.fill.background()
    # bottom stripe
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, Inches(7.32), SLIDE_W, Inches(0.18))
    bot.fill.solid()
    bot.fill.fore_color.rgb = EDC_NAVY
    bot.line.fill.background()
    # page number bottom-left
    add_text(slide, Inches(0.4), Inches(7.0), Inches(2), Inches(0.3),
             f'{page_num} |', size=10, color=MUTED)
    # EDC logo bottom-right
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH,
                                 Inches(12.55), Inches(6.85),
                                 width=Inches(0.65), height=Inches(0.65))


def add_title(slide, bold_part, light_part=None):
    """EDC-style title: '**Bold part** | Lighter continuation'"""
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.5),
                                  Inches(12.2), Inches(1.0))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.1
    r1 = p.add_run()
    r1.text = bold_part
    r1.font.name = 'Helvetica'
    r1.font.size = Pt(34)
    r1.font.bold = True
    r1.font.color.rgb = EDC_NAVY
    if light_part:
        r2 = p.add_run()
        r2.text = f' | {light_part}'
        r2.font.name = 'Helvetica'
        r2.font.size = Pt(34)
        r2.font.bold = False
        r2.font.color.rgb = EDC_NAVY
    return tx


def set_notes(slide, notes):
    nt = slide.notes_slide
    tf = nt.notes_text_frame
    tf.text = notes


# ---------- Slide builders ------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

SLIDES = []


def slide(fn):
    SLIDES.append(fn)
    return fn


# --- 1. Cover (Punkt 1 style — navy, no logo) ---
@slide
def s_cover(s):
    fill_bg(s, EDC_NAVY)
    # "Punkt 1" yellow
    add_text(s, Inches(0), Inches(2.9), SLIDE_W, Inches(0.6),
             'Session 1', size=22, bold=False, color=EDC_YELLOW,
             align=PP_ALIGN.CENTER, font='Menlo')
    # Big title
    add_text(s, Inches(0), Inches(3.6), SLIDE_W, Inches(1.2),
             'Gode vs. dårlige prompts', size=54, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    # Subtitle
    add_text(s, Inches(0), Inches(4.7), SLIDE_W, Inches(0.6),
             'Kickoff og fundamentet · EDC × Claude Code',
             size=20, color=RGBColor(0xC8, 0xD4, 0xE8),
             align=PP_ALIGN.CENTER)
    # Bottom thin line
    bot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, Inches(7.32), SLIDE_W, Inches(0.18))
    bot.fill.solid()
    bot.fill.fore_color.rgb = EDC_YELLOW
    bot.line.fill.background()


# --- 2. Dagsorden ---
@slide
def s_agenda(s):
    fill_bg(s, WHITE)
    add_text(s, Inches(0), Inches(0.7), SLIDE_W, Inches(1.0),
             'Dagsorden', size=46, bold=True, color=EDC_NAVY,
             align=PP_ALIGN.CENTER)
    # Light grey-blue background block for the table
    table_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), Inches(2.0),
                                  Inches(12.13), Inches(4.7))
    table_bg.fill.solid()
    table_bg.fill.fore_color.rgb = EDC_LIGHT_BG
    table_bg.line.fill.background()

    rows = [
        ('0:00–0:05', 'Velkomst', False),
        ('0:05–0:15', 'Foredrag — sådan kører de næste 5 uger + vælg næste præsentant', True),
        ('0:15–0:35', 'Live demo — samme opgave, dårlig vs. god prompt', True),
        ('0:35–0:55', 'Hands-on — solo på din egen maskine', True),
        ('0:55–1:00', 'Del indsigter + take-home', False),
    ]
    row_h = Inches(0.55)
    gap = Inches(0.08)
    start_y = Inches(2.2)
    time_w = Inches(2.0)
    desc_w = Inches(9.6)
    time_x = Inches(0.85)
    desc_x = Inches(2.95)
    for i, (t, w, highlight) in enumerate(rows):
        y = start_y + (row_h + gap) * i
        # time cell
        time_cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, time_x, y, time_w, row_h)
        time_cell.fill.solid()
        time_cell.fill.fore_color.rgb = EDC_NAVY if highlight else EDC_GREY
        time_cell.line.fill.background()
        tt = time_cell.text_frame
        tt.margin_left = Inches(0.2)
        tt.margin_right = Inches(0.2)
        tt.vertical_anchor = MSO_ANCHOR.MIDDLE
        tp = tt.paragraphs[0]
        tp.alignment = PP_ALIGN.LEFT
        tr = tp.add_run()
        tr.text = t
        tr.font.name = 'Helvetica'
        tr.font.size = Pt(15)
        tr.font.bold = True
        tr.font.color.rgb = WHITE
        # description cell
        desc_cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, desc_x, y, desc_w, row_h)
        desc_cell.fill.solid()
        desc_cell.fill.fore_color.rgb = EDC_NAVY if highlight else EDC_GREY
        desc_cell.line.fill.background()
        dt = desc_cell.text_frame
        dt.margin_left = Inches(0.25)
        dt.margin_right = Inches(0.2)
        dt.vertical_anchor = MSO_ANCHOR.MIDDLE
        dp = dt.paragraphs[0]
        dp.alignment = PP_ALIGN.LEFT
        dr = dp.add_run()
        dr.text = w
        dr.font.name = 'Helvetica'
        dr.font.size = Pt(15)
        dr.font.bold = True
        dr.font.color.rgb = WHITE


# --- 3. Vores rolle ---
@slide
def s_role(s):
    fill_bg(s, WHITE)
    add_title(s, 'Vores rolle', 'Facilitatorer, ikke eksperter')
    # Two big content boxes (mirror EDC style page 3)
    box_y_1 = Inches(2.0)
    box_y_2 = Inches(4.4)
    box_w = Inches(12.13)
    box_h = Inches(2.1)
    box_x = Inches(0.6)

    def role_box(y, num, head, body):
        b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, box_x, y, box_w, box_h)
        b.fill.solid()
        b.fill.fore_color.rgb = EDC_LIGHT_BG
        b.line.fill.background()
        # number circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    box_x + Inches(0.4), y + Inches(0.7),
                                    Inches(0.7), Inches(0.7))
        circle.fill.solid()
        circle.fill.fore_color.rgb = EDC_NAVY
        circle.line.fill.background()
        ct = circle.text_frame
        ct.margin_left = 0
        ct.margin_right = 0
        ct.margin_top = 0
        ct.margin_bottom = 0
        ct.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ct.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = num
        cr.font.name = 'Helvetica'
        cr.font.size = Pt(20)
        cr.font.bold = True
        cr.font.color.rgb = WHITE
        # heading + body
        add_text(s, box_x + Inches(1.5), y + Inches(0.35),
                 box_w - Inches(1.8), Inches(0.5),
                 head, size=22, bold=True, color=EDC_NAVY)
        add_rich_text(s, box_x + Inches(1.5), y + Inches(0.95),
                      box_w - Inches(1.8), Inches(1.3),
                      body, size=16, color=INK, line_spacing=1.4)

    role_box(box_y_1, '1',
             'Vi er facilitatorer hele forløbet',
             'Michael og jeg står for **selve faciliteringen** alle 5 uger. Vi er her for at gøre det trygt og holde tråden — det her er et hold der lærer sammen.')
    role_box(box_y_2, '2',
             'I underviser hinanden fra session 2',
             'Fra session 2 er det **én deltager fra holdet** der præsenterer dagens emne. Vi hjælper med slides ugen før, så ingen står alene med det. **Pointe:** I lærer det bedst ved at undervise det.')


# --- 4. Vælg næste præsentant ---
@slide
def s_next(s):
    fill_bg(s, WHITE)
    add_title(s, 'Nu', 'Vælg næste præsentant')
    # Big highlighted box
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.6), Inches(2.2),
                             Inches(12.13), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = EDC_LIGHT_BG
    box.line.fill.background()
    # Star icon (yellow)
    star = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,
                              Inches(1.0), Inches(2.6),
                              Inches(1.1), Inches(1.1))
    star.fill.solid()
    star.fill.fore_color.rgb = EDC_YELLOW
    star.line.fill.background()

    add_text(s, Inches(2.6), Inches(2.5),
             Inches(10), Inches(0.6),
             'Hvem vil tage session 2?', size=14, bold=True,
             color=EDC_YELLOW, font='Menlo')
    add_text(s, Inches(2.6), Inches(2.95),
             Inches(10), Inches(0.7),
             'Plan mode vs. ikke plan mode',
             size=28, bold=True, color=EDC_NAVY)
    add_text(s, Inches(2.6), Inches(3.55),
             Inches(10), Inches(0.5),
             'Torsdag · uge 1',
             size=18, color=MUTED)
    add_rich_text(s, Inches(2.6), Inches(4.3),
                  Inches(10), Inches(1.2),
                  'Jeg sender slide-skabelon og pre-work-video **senest dagen før**, så du står ikke alene med det. Hvis ingen melder sig, vælger vi én — og vi gentager øvelsen i starten af hver session.',
                  size=15, color=INK, line_spacing=1.5)


# --- 5. Dagens kerne ---
@slide
def s_core(s):
    fill_bg(s, WHITE)
    add_title(s, 'Dagens kerne', 'En prompt er alt')
    # Two big lines
    add_text(s, Inches(0.6), Inches(2.7),
             Inches(12.2), Inches(0.9),
             'En dårlig prompt giver generisk kode.',
             size=34, bold=True, color=MUTED)
    add_text(s, Inches(0.6), Inches(3.7),
             Inches(12.2), Inches(0.9),
             'En god prompt giver produktionsklar kode.',
             size=34, bold=True, color=EDC_NAVY)
    # Subtitle
    sub = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.6), Inches(5.5),
                             Inches(12.13), Inches(1.0))
    sub.fill.solid()
    sub.fill.fore_color.rgb = EDC_LIGHT_BG
    sub.line.fill.background()
    add_text(s, Inches(0.6), Inches(5.5), Inches(12.13), Inches(1.0),
             'Forskellen er hvor præcist du beskriver opgaven.',
             size=22, color=EDC_NAVY, align=PP_ALIGN.CENTER,
             anchor=MSO_ANCHOR.MIDDLE)


# --- 6. De 4 byggeklodser ---
@slide
def s_blocks(s):
    fill_bg(s, WHITE)
    add_title(s, 'De 4 byggeklodser', 'Tænk i fire dele når du skriver en prompt')
    blocks = [
        ('1', 'Kontekst', None,
         'Hvilke filer skal\nClaude læse?'),
        ('2', 'Opgave', None,
         'Hvad vil du have\ngjort? Fx "Lav en\nhelper", "Tilføj knap".'),
        ('3', 'Begrænsninger', 'Når relevant',
         'Hvad må Claude\nIKKE røre?'),
        ('4', 'Forventet output', None,
         'Hvad ser "done"\nud som?'),
    ]
    box_w = Inches(2.95)
    box_h = Inches(3.6)
    gap = Inches(0.13)
    start_x = Inches(0.65)
    top = Inches(2.2)
    for i, (num, head, badge, body) in enumerate(blocks):
        x = start_x + (box_w + gap) * i
        b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, top, box_w, box_h)
        b.fill.solid()
        b.fill.fore_color.rgb = EDC_LIGHT_BG
        b.line.fill.background()
        # number circle
        circle = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                    x + Inches(0.4), top + Inches(0.4),
                                    Inches(0.6), Inches(0.6))
        circle.fill.solid()
        circle.fill.fore_color.rgb = EDC_NAVY
        circle.line.fill.background()
        ct = circle.text_frame
        ct.margin_left = 0
        ct.margin_right = 0
        ct.margin_top = 0
        ct.margin_bottom = 0
        ct.vertical_anchor = MSO_ANCHOR.MIDDLE
        cp = ct.paragraphs[0]
        cp.alignment = PP_ALIGN.CENTER
        cr = cp.add_run()
        cr.text = num
        cr.font.name = 'Helvetica'
        cr.font.size = Pt(20)
        cr.font.bold = True
        cr.font.color.rgb = WHITE
        # optional "Når relevant" badge top-right
        if badge:
            badge_w = Inches(1.2)
            badge_h = Inches(0.3)
            bx = x + box_w - badge_w - Inches(0.25)
            by = top + Inches(0.55)
            badge_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                           bx, by, badge_w, badge_h)
            badge_box.adjustments[0] = 0.5
            badge_box.fill.solid()
            badge_box.fill.fore_color.rgb = EDC_YELLOW
            badge_box.line.fill.background()
            bt = badge_box.text_frame
            bt.margin_left = 0
            bt.margin_right = 0
            bt.margin_top = 0
            bt.margin_bottom = 0
            bt.vertical_anchor = MSO_ANCHOR.MIDDLE
            bp = bt.paragraphs[0]
            bp.alignment = PP_ALIGN.CENTER
            br = bp.add_run()
            br.text = badge
            br.font.name = 'Helvetica'
            br.font.size = Pt(11)
            br.font.bold = True
            br.font.color.rgb = EDC_NAVY
        # heading (smaller font so "Begrænsninger" fits one line)
        add_text(s, x + Inches(0.3), top + Inches(1.3),
                 box_w - Inches(0.6), Inches(0.7),
                 head, size=18, bold=True, color=EDC_NAVY)
        # body
        add_text(s, x + Inches(0.3), top + Inches(2.0),
                 box_w - Inches(0.6), Inches(1.4),
                 body, size=14, color=INK, line_spacing=1.4)
    # mantra
    add_text(s, Inches(0.6), Inches(6.15),
             Inches(12.13), Inches(0.6),
             '[Kontekst] + [Opgave] + [Begrænsninger] + [Forventet output]',
             size=15, color=MUTED, align=PP_ALIGN.CENTER, font='Menlo')


# --- 7. Tryghed: Esc Esc ---
@slide
def s_esc(s):
    fill_bg(s, WHITE)
    add_title(s, 'Tryghed fra start', 'Esc · Esc — undo for hele sessionen')
    # Big code-style label box
    keybox = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                Inches(0.7), Inches(2.2),
                                Inches(4.5), Inches(2.5))
    keybox.adjustments[0] = 0.08
    keybox.fill.solid()
    keybox.fill.fore_color.rgb = EDC_NAVY
    keybox.line.fill.background()
    add_text(s, Inches(0.7), Inches(2.2), Inches(4.5), Inches(2.5),
             'Esc · Esc', size=64, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font='Menlo')
    # Bullets
    add_bullets(s, Inches(5.6), Inches(2.3), Inches(7.2), Inches(2.5), [
        'Tryk **Esc** to gange → rewind-menu',
        'Vælg et tidligere step → filerne ruller tilbage',
        'Vi går i dybden i **session 3** — men brug det fra dag 1',
    ], size=18, line_spacing=1.5)
    # bottom note
    note = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(5.2),
                              Inches(12.13), Inches(1.3))
    note.fill.solid()
    note.fill.fore_color.rgb = EDC_LIGHT_BG
    note.line.fill.background()
    add_rich_text(s, Inches(1.0), Inches(5.2),
                  Inches(11.5), Inches(1.3),
                  'Det er det her der gør det **trygt** at lade Claude prøve noget vildt. Hvis I aldrig kan rulle tilbage, tør I aldrig prøve noget nyt.',
                  size=18, color=EDC_NAVY, anchor=MSO_ANCHOR.MIDDLE,
                  line_spacing=1.4)


# --- 8. Section divider: Live demo ---
@slide
def s_divider_demo(s):
    fill_bg(s, EDC_NAVY)
    if os.path.exists(LOGO_PATH):
        logo_w = Inches(1.6)
        s.shapes.add_picture(LOGO_PATH,
                             (SLIDE_W - logo_w) / 2, Inches(1.8),
                             width=logo_w, height=logo_w)
    add_text(s, Inches(0), Inches(4.0), SLIDE_W, Inches(0.6),
             'Punkt 2', size=22, color=EDC_YELLOW,
             align=PP_ALIGN.CENTER, font='Menlo')
    add_text(s, Inches(0), Inches(4.7), SLIDE_W, Inches(1.0),
             'Live demo', size=44, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(5.5), SLIDE_W, Inches(0.6),
             'Samme opgave — dårlig prompt vs. god prompt',
             size=18, color=RGBColor(0xC8, 0xD4, 0xE8),
             align=PP_ALIGN.CENTER)
    bot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, Inches(7.32), SLIDE_W, Inches(0.18))
    bot.fill.solid()
    bot.fill.fore_color.rgb = EDC_YELLOW
    bot.line.fill.background()


# --- 9. Bug'en på /demo ---
@slide
def s_bug(s):
    fill_bg(s, WHITE)
    add_title(s, "Bug'en på /demo", 'Priser i USD med cents på en dansk side')
    col_w = Inches(5.85)
    left_x = Inches(0.7)
    right_x = Inches(6.85)
    top_y = Inches(2.3)
    box_h = Inches(3.7)
    # NU
    bug = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, top_y, col_w, box_h)
    bug.fill.solid()
    bug.fill.fore_color.rgb = RED_BG
    bug.line.fill.background()
    add_text(s, left_x + Inches(0.4), top_y + Inches(0.3),
             col_w - Inches(0.8), Inches(0.4),
             'NU (BUG)', size=14, bold=True, color=RED_INK)
    add_text(s, left_x + Inches(0.4), top_y + Inches(0.95),
             col_w - Inches(0.8), Inches(2.5),
             '$8,500,000.00\n$4,250,000.00\n$6,900,000.00',
             size=30, bold=True, color=RED_INK, font='Menlo',
             line_spacing=1.3)
    # ØNSKET
    ok = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, top_y, col_w, box_h)
    ok.fill.solid()
    ok.fill.fore_color.rgb = GREEN_BG
    ok.line.fill.background()
    add_text(s, right_x + Inches(0.4), top_y + Inches(0.3),
             col_w - Inches(0.8), Inches(0.4),
             'ØNSKET', size=14, bold=True, color=GREEN_INK)
    add_text(s, right_x + Inches(0.4), top_y + Inches(0.95),
             col_w - Inches(0.8), Inches(2.5),
             '8.500.000 kr.\n4.250.000 kr.\n6.900.000 kr.',
             size=30, bold=True, color=GREEN_INK, font='Menlo',
             line_spacing=1.3)
    # Note below
    add_text(s, Inches(0.7), Inches(6.25),
             Inches(12.0), Inches(0.5),
             'Pris-formateringen ligger inline 3 steder uden helper — typisk consequence af "hurtigt fix".',
             size=14, italic=True, color=MUTED,
             align=PP_ALIGN.CENTER)


# --- 10. Dårlig prompt ---
@slide
def s_bad_prompt(s):
    fill_bg(s, WHITE)
    add_title(s, 'Forsøg 1', 'Dårlig prompt — "realistisk dårlig"')
    add_code_block(s, Inches(0.7), Inches(2.2), Inches(11.93), Inches(2.0),
                   'Priserne på /demo vises som amerikanske dollars med cents\n'
                   '(fx "$8,500,000.00") overalt. Det er en dansk side.\n'
                   'Fix det. Ingen commit eller push.',
                   size=18)
    # Hvad sker der
    add_text(s, Inches(0.7), Inches(4.5), Inches(12), Inches(0.5),
             'Hvad sker der?', size=20, bold=True, color=EDC_NAVY)
    add_bullets(s, Inches(0.7), Inches(5.0), Inches(12.0), Inches(2.0), [
        'Claude finder *et* fix — men forskelligt hver gang du kører prompten',
        'Inline-duplikering 3 steder, eller patch kun én fil og misser resten',
        'Det "virker" — men efterlader kode der er svær at vedligeholde',
    ], size=16, line_spacing=1.5)


# --- 11. God prompt ---
@slide
def s_good_prompt(s):
    fill_bg(s, WHITE)
    add_title(s, 'Forsøg 2', 'God prompt — samme opgave, med kontekst')
    add_code_block(s, Inches(0.7), Inches(1.9), Inches(11.93), Inches(4.4),
                   "Pris-formateringen ligger inline tre steder med\n"
                   "Intl.NumberFormat('en-US', { style: 'currency',\n"
                   "currency: 'USD', maximumFractionDigits: 2 }):\n"
                   '\n'
                   '- lib/propertyService.ts (i getAll, formattedPrice)\n'
                   '- components/PropertyCard.tsx\n'
                   '- app/properties/[id]/page.tsx (formattedPrice)\n'
                   '\n'
                   "Lav en formatPrice(price: number)-helper i\n"
                   "lib/propertyService.ts der bruger 'da-DK', 'DKK' og\n"
                   "maximumFractionDigits: 0. Brug den fra alle tre steder.\n"
                   '\n'
                   'Rør intet andet. Ingen commit eller push.',
                   size=14)
    # Result
    note = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.7), Inches(6.5),
                              Inches(11.93), Inches(0.65))
    note.fill.solid()
    note.fill.fore_color.rgb = EDC_LIGHT_BG
    note.line.fill.background()
    add_text(s, Inches(0.7), Inches(6.5), Inches(11.93), Inches(0.65),
             'Resultat: én helper, brugt 3 steder. Deterministisk.',
             size=16, bold=True, color=EDC_NAVY,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# --- 12. Pointen ---
@slide
def s_point(s):
    fill_bg(s, WHITE)
    add_title(s, 'Pointen', 'Begge fixes virker visuelt — kun den ene efterlader sund kode')
    col_w = Inches(5.85)
    left_x = Inches(0.7)
    right_x = Inches(6.85)
    top_y = Inches(2.2)
    box_h = Inches(4.0)
    # bad
    bad = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, top_y, col_w, box_h)
    bad.fill.solid()
    bad.fill.fore_color.rgb = RED_BG
    bad.line.fill.background()
    add_text(s, left_x + Inches(0.4), top_y + Inches(0.3),
             col_w - Inches(0.8), Inches(0.4),
             'DÅRLIG PROMPT', size=14, bold=True, color=RED_INK)
    add_text(s, left_x + Inches(0.4), top_y + Inches(0.85),
             col_w - Inches(0.8), Inches(0.9),
             '3 inline-kopier af samme logik',
             size=22, bold=True, color=EDC_NAVY)
    add_bullets(s, left_x + Inches(0.4), top_y + Inches(2.0),
                col_w - Inches(0.8), Inches(2.0), [
        'DRY brudt — ændring kræver 3 filer',
        'Forskelligt resultat hver kørsel',
        'Nogle gange ufuldstændig fix',
    ], size=15, line_spacing=1.5, bullet_color=RED_INK)
    # good
    good = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, top_y, col_w, box_h)
    good.fill.solid()
    good.fill.fore_color.rgb = GREEN_BG
    good.line.fill.background()
    add_text(s, right_x + Inches(0.4), top_y + Inches(0.3),
             col_w - Inches(0.8), Inches(0.4),
             'GOD PROMPT', size=14, bold=True, color=GREEN_INK)
    add_text(s, right_x + Inches(0.4), top_y + Inches(0.85),
             col_w - Inches(0.8), Inches(0.9),
             '1 helper, brugt 3 steder',
             size=22, bold=True, color=EDC_NAVY)
    add_bullets(s, right_x + Inches(0.4), top_y + Inches(2.0),
                col_w - Inches(0.8), Inches(2.0), [
        'Ét sted at ændre fremover',
        'Deterministisk — samme resultat hver gang',
        'Helperen kan genbruges',
    ], size=15, line_spacing=1.5, bullet_color=GREEN_INK)


# --- 13. Vigtigt forbehold ---
@slide
def s_caveat(s):
    fill_bg(s, WHITE)
    add_title(s, 'Vigtigt forbehold', 'Demoen er lille — virkeligheden er stor')
    # Star icon + intro
    star = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,
                              Inches(0.8), Inches(2.1),
                              Inches(0.9), Inches(0.9))
    star.fill.solid()
    star.fill.fore_color.rgb = EDC_NAVY
    star.line.fill.background()
    add_rich_text(s, Inches(2.1), Inches(2.15),
                  Inches(10.6), Inches(0.9),
                  '**Forskellen ser subtil ud her.** På jeres rigtige projekter eksploderer den.',
                  size=20, color=EDC_NAVY, line_spacing=1.3)

    # Bullets
    add_bullets(s, Inches(0.8), Inches(3.5), Inches(11.8), Inches(3.5), [
        '**Her:** 5 filer, ét tydeligt symptom, Opus 4.7 — Claude finder *et* fix uanset hvor vag prompten er',
        '**På jeres codebases:** 100k+ linjer, 10+ års historik, DTO\'er, mappere, migrations, uskrevne konventioner',
        '**Konsekvens af vag prompt:** forkert fil · glemt halvdel af kæden · konfliktende abstraktion · brudt konvention',
        '**Vi træner mønsteret HER** — så det sidder i fingrene mandag morgen i en rigtig PR',
    ], size=15, line_spacing=1.6)


# --- 14. Section divider: Hands-on ---
@slide
def s_divider_handson(s):
    fill_bg(s, EDC_NAVY)
    if os.path.exists(LOGO_PATH):
        logo_w = Inches(1.6)
        s.shapes.add_picture(LOGO_PATH,
                             (SLIDE_W - logo_w) / 2, Inches(1.8),
                             width=logo_w, height=logo_w)
    add_text(s, Inches(0), Inches(4.0), SLIDE_W, Inches(0.6),
             'Punkt 3', size=22, color=EDC_YELLOW,
             align=PP_ALIGN.CENTER, font='Menlo')
    add_text(s, Inches(0), Inches(4.7), SLIDE_W, Inches(1.0),
             'Hands-on', size=44, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(5.5), SLIDE_W, Inches(0.6),
             '20 minutter · solo på din egen maskine',
             size=18, color=RGBColor(0xC8, 0xD4, 0xE8),
             align=PP_ALIGN.CENTER)
    bot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, Inches(7.32), SLIDE_W, Inches(0.18))
    bot.fill.solid()
    bot.fill.fore_color.rgb = EDC_YELLOW
    bot.line.fill.background()


# --- 15. Hands-on instructions ---
@slide
def s_handson(s):
    fill_bg(s, WHITE)
    add_title(s, 'Hands-on', 'Prøv begge prompt-typer på samme opgave')
    add_bullets(s, Inches(0.7), Inches(2.2), Inches(12.0), Inches(4.5), [
        '**git clone** + **npm install** + **claude** i `demo/`-mappen',
        'Vælg én opgave fra demoen (pris-bug, performance, eller favorit-knap)',
        'Skriv først en **hurtig prompt** — som du ville skrive på en travl dag',
        '**git checkout .** og **/clear** for at rulle tilbage',
        'Skriv så en prompt med **alle 4 byggeklodser**',
        'Sammenlign: hvad var forskellen? Hvilken ville du committe?',
    ], size=18, line_spacing=1.6)


# --- 16. Hjemmeopgave + tak ---
@slide
def s_outro(s):
    fill_bg(s, EDC_NAVY)
    if os.path.exists(LOGO_PATH):
        logo_w = Inches(1.2)
        s.shapes.add_picture(LOGO_PATH,
                             (SLIDE_W - logo_w) / 2, Inches(0.7),
                             width=logo_w, height=logo_w)
    add_text(s, Inches(0), Inches(2.1), SLIDE_W, Inches(0.6),
             'Hjemmeopgave inden torsdag', size=18, color=EDC_YELLOW,
             align=PP_ALIGN.CENTER, font='Menlo')
    add_text(s, Inches(0), Inches(2.7), SLIDE_W, Inches(1.0),
             'Brug Claude Code', size=44, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(3.5), SLIDE_W, Inches(0.7),
             'på mindst én rigtig opgave', size=28, color=RGBColor(0xC8, 0xD4, 0xE8),
             align=PP_ALIGN.CENTER)
    # Box with details
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(2.0), Inches(4.7),
                             Inches(9.33), Inches(2.0))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1B, 0x36, 0x6B)
    box.line.fill.background()
    add_text(s, Inches(2.4), Inches(4.85),
             Inches(8.5), Inches(0.5),
             'Tag med til torsdag:', size=16, bold=True, color=EDC_YELLOW,
             font='Menlo')
    add_bullets(s, Inches(2.4), Inches(5.4),
                Inches(8.5), Inches(1.3), [
        'Én prompt der virkede godt',
        'Én der ikke gjorde — vi lærer mest af de fejlede',
    ], size=18, color=WHITE, line_spacing=1.4,
                bullet_color=EDC_YELLOW)
    bot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, Inches(7.32), SLIDE_W, Inches(0.18))
    bot.fill.solid()
    bot.fill.fore_color.rgb = EDC_YELLOW
    bot.line.fill.background()


# ---------- Speaker notes -------------------------------------------------
# Format-konvention:
#   ─── LÆS ─── = baggrund/kontekst du skanner med øjnene
#   ─── SIG ─── = ordret hvad du siger (i citationstegn)
#   ─── GØR ─── = handlinger (klik, kopier, refresh, osv.)
NOTES = [
    # 1. Cover
    "─── LÆS ───\n"
    "Velkomst · 0:00–0:05. Første øjeblik — bare rammesæt. Vent ikke for længe; kom videre.\n\n"
    "─── SIG ───\n"
    "\"Velkommen. Jeg er Matias, og det her er Michael. Vi er facilitatorer for de næste 5 uger.\"\n\n"
    "\"Det her er ikke et kursus hvor I lytter til os. Det er en træningslejr hvor I lærer ved at gøre — og fra session 2 også ved at undervise jeres kollegaer.\"",

    # 2. Dagsorden
    "─── LÆS ───\n"
    "Hurtigt overblik · 0:05–0:07. Peg på de fremhævede (blå) felter — det er de tre hoveddele. Ingen pause på en time, men det går stærkt.\n\n"
    "─── SIG ───\n"
    "\"Lige om lidt fortæller jeg lidt om hvordan de næste 5 uger ser ud.\"\n\n"
    "\"Så går vi i terminalen og laver en live demo.\"\n\n"
    "\"Og så er det jer der skal kode — solo på jeres egen maskine.\"",

    # 3. Vores rolle
    "─── LÆS ───\n"
    "Vigtigste pointe i dag · 0:07–0:10. Hvis de bare lytter, lærer de ingenting. Hvis de underviser hinanden, sidder det fast.\n\n"
    "─── SIG ───\n"
    "\"Det her er den vigtigste pointe i dag: Michael og jeg er facilitatorer, ikke eksperter.\"\n\n"
    "\"Jeg ved ikke nødvendigvis mere end jer om Claude Code om 5 uger. Det her er et hold der lærer sammen.\"\n\n"
    "\"Fra session 2 er det én af jer der præsenterer dagens emne. Vi hjælper med slides ugen før, så ingen står alene med det.\"",

    # 4. Vælg næste præsentant
    "─── LÆS ───\n"
    "Akavet øjeblik · 0:10–0:13. Stil spørgsmålet og hold mund. Lad stilheden virke. Hvis ingen melder sig efter 10 sek: vælg én du ved godt har lyst — eller én der ser frisk ud.\n\n"
    "─── SIG ───\n"
    "\"Hvem vil tage session 2 om Plan mode? Det er torsdag, jeg sender slide-skabelon og pre-work-video senest dagen før.\"\n\n"
    "[Vent. Tæl til 10 i hovedet.]\n\n"
    "Hvis ingen melder sig:\n"
    "\"[Navn] — vil du tage den?\"\n\n"
    "─── GØR ───\n"
    "Når én melder sig: TAK højt. Det sætter normen for resten af forløbet.",

    # 5. Dagens kerne
    "─── LÆS ───\n"
    "Saml i ÉN sætning · 0:13–0:15. Sig den langsomt. Det er hele dagen i én linje.\n\n"
    "─── SIG ───\n"
    "\"En dårlig prompt giver generisk kode.\"\n\n"
    "\"En god prompt giver produktionsklar kode.\"\n\n"
    "\"Forskellen er hvor præcist du beskriver opgaven for Claude. Det er det vi træner i dag.\"",

    # 6. De 4 byggeklodser
    "─── LÆS ───\n"
    "Mønsteret · 0:15–0:18. Det her skal de kunne udenad efter session 1. Læs dem op én af gangen og giv et lille eksempel for hver. Vigtigt: BEGRÆNSNINGER er markeret \"når relevant\" — det er ikke nødvendigt for hver prompt (fx ikke for \"forklar denne fil\"), men ved kode-ændringer er det guld værd.\n\n"
    "─── SIG ───\n"
    "\"Tænk i fire dele når du skriver en prompt:\"\n\n"
    "\"1. Kontekst — fx: 'Læs lib/propertyService.ts og components/PropertyCard.tsx'.\"\n\n"
    "\"2. Opgave — konkret, fx: 'Lav en formatPrice-helper og brug den i begge filer'. Ikke 'der er noget med priserne...'.\"\n\n"
    "\"3. Begrænsninger — når relevant. Fx: 'Rør intet andet. Ingen commit.' Hvis du bare vil have Claude til at forklare en fil, kan du droppe det her.\"\n\n"
    "\"4. Forventet output — fx: 'Priser vises som 8.500.000 kr. alle 3 steder.'\"\n\n"
    "\"Vi vender tilbage til mønsteret hver eneste session de næste 5 uger.\"",

    # 7. Esc Esc
    "─── LÆS ───\n"
    "Tryghed fra dag 1 · 0:18–0:20. Vi går i dybden i session 3, men de skal vide om det FRA i dag. Pointe: hvis de aldrig kan rulle tilbage, tør de aldrig prøve noget nyt.\n\n"
    "─── SIG ───\n"
    "\"Når Claude laver noget I fortryder — eller koden går i sort — så tryk Esc to gange.\"\n\n"
    "\"Der kommer en menu hvor I kan rulle tilbage til ethvert tidligere punkt i samtalen. Filerne, conversation, eller begge dele.\"\n\n"
    "\"Det er det her der gør det trygt at lade Claude prøve noget vildt.\"\n\n"
    "\"I må gerne bruge det fra første prompt I skriver i dag.\"",

    # 8. Live demo divider
    "─── LÆS ───\n"
    "Skift til terminal · 0:20 — start.\n\n"
    "TJEKLISTE FØR DU KØRER:\n"
    "  □ git status er clean\n"
    "  □ /demo er åben i browseren med USD-priser synlige\n"
    "  □ claude --dangerously-skip-permissions kører i demo/-mappen\n"
    "  □ Du kan se BÅDE terminal og browser samtidig på projektoren\n\n"
    "─── SIG ───\n"
    "\"OK — nu skifter vi til terminal og browser.\"\n\n"
    "\"I de næste 15 minutter skriver jeg samme opgave til Claude — først dårligt, så godt.\"\n\n"
    "\"Kig BÅDE på terminalen og på /demo i browseren ved siden af.\"",

    # 9. Bug'en
    "─── LÆS ───\n"
    "Vis bug'en · 0:20–0:22.\n\n"
    "─── GØR ───\n"
    "Skift til browseren. Scroll ned ad /demo. Peg eksplicit på priserne.\n"
    "Klik på en bolig — vis at detaljesiden har samme problem.\n"
    "Søg \"demo@edc.dk\" i favorit-widgeten — også USD-priser.\n\n"
    "─── SIG ───\n"
    "\"Det her er en dansk bolig-side. Priserne er i amerikanske dollars MED cents. Hele sitet er i stykker visuelt.\"\n\n"
    "\"Pris-formateringen ligger inline 3 steder uden helper. Det er ikke konstrueret — det er typisk consequence af 'hurtigt fix' nogen lavede engang.\"\n\n"
    "\"Det er det her vi vil have Claude til at fixe.\"",

    # 10. Dårlig prompt
    "─── LÆS ───\n"
    "Realistisk dårlig · 0:22–0:26. Vi har kørt det her 3 gange — den lander forskelligt hver gang.\n\n"
    "─── SIG ───\n"
    "\"Det her ligner en prompt de fleste skriver på autopilot — vi nævner symptomet, men ikke hvor i koden det skal fixes.\"\n\n"
    "\"Det er IKKE en stråmand, det er hverdag.\"\n\n"
    "─── GØR ───\n"
    "Kopier prompten fra slidet ind i Claude. Tryk enter.\n\n"
    "─── SIG (mens Claude arbejder) ───\n"
    "\"Forventning: den finder ET fix der virker. Men hold den her i baghovedet — vi har kørt det 3 gange, den lander forskelligt hver gang.\"\n\n"
    "─── GØR ───\n"
    "Når Claude er færdig: vis hvad den gjorde, refresh /demo.\n\n"
    "─── SIG ───\n"
    "\"Det virker! Men hold det her i baghovedet.\"",

    # 11. God prompt
    "─── LÆS ───\n"
    "Samme opgave, med kontekst · 0:26–0:31.\n\n"
    "─── GØR ───\n"
    "Rul tilbage i terminalen (vis det højt mens du gør det):\n"
    "  git checkout .\n"
    "  /clear\n\n"
    "─── SIG ───\n"
    "\"Nu siger vi præcis HVOR (3 filer), HVAD (én helper), og HVORDAN (da-DK, DKK, 0 decimaler).\"\n\n"
    "─── GØR ───\n"
    "Kopier den gode prompt ind. Tryk enter.\n"
    "Når den er færdig: vis at den lavede ÉN helper i propertyService.ts, og at de andre 2 filer importerer den. Refresh /demo.\n\n"
    "─── SIG ───\n"
    "\"Samme visuelle resultat. Helt anden kode-struktur.\"",

    # 12. Pointen
    "─── LÆS ───\n"
    "Sig det højt · 0:31–0:33. Det her er momentet hvor du SKAL sige det højt — ikke bare lade slidet tale. Det er hele dagen i én sætning.\n\n"
    "─── SIG ───\n"
    "\"Begge fixes virker visuelt — kr. står overalt. Men kig på koden:\"\n\n"
    "\"Den dårlige duplikerede logikken 3 steder.\"\n\n"
    "\"Den gode lavede én helper, brugt 3 steder.\"\n\n"
    "\"Om 3 måneder, når en designer beder om at vise øre i favorit-listen, er det den forskel der bestemmer om I bruger 5 minutter eller 30.\"",

    # 13. Vigtigt forbehold
    "─── LÆS ───\n"
    "Forsikring · 0:33–0:35. Det her slide er forsikring mod den deltager der siger \"jamen den dårlige fungerede jo fint\". Anerkend det direkte — det er ærligt.\n\n"
    "─── SIG ───\n"
    "\"I har ret. På denne lille demo med Opus 4.7 er forskellen subtil — Claude er smart nok til at finde et fix der virker, selv når prompten er vag.\"\n\n"
    "\"Men i jeres rigtige projekter — 100k+ linjer, DTO'er, mappere, migrations, konventioner der ikke er skrevet ned — eksploderer det.\"\n\n"
    "\"En vag prompt rammer den forkerte fil, glemmer halvdelen af kæden, eller introducerer et abstraktionslag der konflikter med jeres eksisterende.\"\n\n"
    "\"Vi træner mønsteret her hvor det er sikkert at fejle — så det sidder i fingrene når I er midt i en rigtig PR mandag morgen.\"",

    # 14. Hands-on divider
    "─── LÆS ───\n"
    "Korte instrukser · 0:35 — start. De skal i gang.\n\n"
    "─── SIG ───\n"
    "\"Nu jer.\"\n\n"
    "\"Solo på din egen maskine — du skal selv mærke forskellen i fingrene.\"\n\n"
    "\"I har 20 minutter.\"",

    # 15. Hands-on instructions
    "─── LÆS ───\n"
    "Mens de arbejder · 0:35–0:55. Gå rundt. Hjælp dem der sidder fast. Lyt efter en god 'aha'-bemærkning du kan bede dem dele bagefter.\n\n"
    "─── SIG (i starten) ───\n"
    "\"Læs trinene på slidet: clone, npm install, claude i demo. Vælg én opgave fra demoen.\"\n\n"
    "\"Det vigtigste: prøv BEGGE prompts — ikke bare den gode. Læg mærke til forskellen.\"\n\n"
    "\"Hvis du sidder fast, ræk hånden op — vi går rundt.\"\n\n"
    "─── SIG (med 5 min tilbage) ───\n"
    "\"Saml jer — én indsigt højt fra hver. Den der overraskede jer mest.\"",

    # 16. Outro
    "─── LÆS ───\n"
    "Afslutning · 0:55–1:00. Hjemmeopgave nem. Bekræft næste præsentant igen ved navn så de ved det er reelt.\n\n"
    "─── SIG ───\n"
    "\"Hjemmeopgaven er nem. Brug Claude Code på MINDST én rigtig opgave inden torsdag — i et af jeres rigtige projekter, ikke demoen.\"\n\n"
    "\"Tag to prompts med til torsdag: én der virkede godt, og én der ikke gjorde.\"\n\n"
    "\"Vi starter session 2 med at dele dem. Vi lærer mest af de fejlede prompts.\"\n\n"
    "\"[Navn] tager session 2 om Plan mode. Jeg sender slide-skabelon og pre-work-video senest mandag.\"\n\n"
    "\"Tak for i dag.\"",
]


# ---------- Build deck ----------------------------------------------------
total = len(SLIDES)
content_pages = {0, 7, 13, total - 1}  # cover + dividers + outro have own chrome

for idx, builder in enumerate(SLIDES):
    s = prs.slides.add_slide(blank)
    builder(s)
    page_num = idx + 1
    # Add chrome only to standard content slides
    if idx not in {0, 7, 13, total - 1}:
        add_chrome(s, page_num)
    if idx < len(NOTES):
        set_notes(s, NOTES[idx])

out = os.path.join(HERE, 'session1.pptx')
prs.save(out)
print(f'Wrote {out}  ({total} slides)')
