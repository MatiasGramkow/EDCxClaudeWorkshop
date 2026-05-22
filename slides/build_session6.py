"""Generate session 6 PowerPoint deck for the EDC × Claude Code workshop.

Session 6: Git, commits, review og PAT (Troels Ford Helbo præsenterer).
Bruger samme EDC-design som session 1-5.

Troels' brief er to ben + live demoer til begge:
  1. PAT mod Azure DevOps + hvordan Claude læser DevOps work items
  2. Hvordan Claude bruger git med branches, commits og reviews

Troels får hele decket stukket. Hans content-slides (6-13) og hands-on
cheat-sheet (16) er skjulte fordi han bygger sine egne — slides ligger
som backup hvis han skal bruge dem.

Hands-on knytter sig til workshop-lab'en: deltagerne pusher deres
første PR til /demo/lab/<deres-slug>.

Run: python3 build_session6.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import os

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH = os.path.join(HERE, 'edc_logo.png')

# ---------- EDC palette --------------------------------------------------
EDC_NAVY = RGBColor(0x12, 0x2A, 0x52)
EDC_NAVY_DARK = RGBColor(0x0A, 0x1F, 0x3D)
EDC_BLUE_BODY = RGBColor(0x1F, 0x3A, 0x6B)
EDC_YELLOW = RGBColor(0xF5, 0xB8, 0x1C)
EDC_LIGHT_BG = RGBColor(0xEE, 0xF1, 0xF6)
EDC_GREY = RGBColor(0x9C, 0xA3, 0xAF)
INK = RGBColor(0x12, 0x2A, 0x52)
MUTED = RGBColor(0x55, 0x6B, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN_BG = RGBColor(0xD7, 0xEC, 0xDD)
GREEN_INK = RGBColor(0x14, 0x53, 0x2D)
RED_BG = RGBColor(0xFA, 0xDC, 0xDC)
RED_INK = RGBColor(0x8B, 0x1F, 0x1F)
CODE_BG = RGBColor(0x0F, 0x1F, 0x3D)
CODE_INK = RGBColor(0xE2, 0xE8, 0xF0)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ---------- Low-level helpers --------------------------------------------
def fill_bg(slide, color):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    return bg


def add_text(slide, left, top, width, height, text, *,
             font='Helvetica', size=24, bold=False, italic=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.15):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tx


def _parse_bold(text):
    out = []
    i = 0
    while i < len(text):
        if text[i:i+2] == '**':
            end = text.find('**', i + 2)
            if end == -1:
                out.append((text[i:], False))
                break
            out.append((text[i+2:end], True))
            i = end + 2
        else:
            nxt = text.find('**', i)
            if nxt == -1:
                out.append((text[i:], False))
                break
            out.append((text[i:nxt], False))
            i = nxt
    return [seg for seg in out if seg[0]]


def add_rich_text(slide, left, top, width, height, text, *,
                  font='Helvetica', size=22, color=INK,
                  align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                  line_spacing=1.3):
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    for seg, is_bold in _parse_bold(text):
        run = p.add_run()
        run.text = seg
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = is_bold
        run.font.color.rgb = color
    return tx


def add_bullets(slide, left, top, width, height, bullets, *,
                font='Helvetica', size=20, color=INK, line_spacing=1.4,
                bullet_glyph='■', bullet_color=None):
    if bullet_color is None:
        bullet_color = EDC_NAVY
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(10)
        br = p.add_run()
        br.text = f'{bullet_glyph}   '
        br.font.name = font
        br.font.size = Pt(size)
        br.font.bold = True
        br.font.color.rgb = bullet_color
        for seg, is_bold in _parse_bold(b):
            run = p.add_run()
            run.text = seg
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = is_bold
            run.font.color.rgb = color
    return tx


def add_code_block(slide, left, top, width, height, code, *, size=14):
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.adjustments[0] = 0.04
    box.fill.solid()
    box.fill.fore_color.rgb = CODE_BG
    box.line.fill.background()
    tf = box.text_frame
    tf.margin_left = Inches(0.4)
    tf.margin_right = Inches(0.4)
    tf.margin_top = Inches(0.3)
    tf.margin_bottom = Inches(0.3)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    lines = code.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.25
        run = p.add_run()
        run.text = line if line else ' '
        run.font.name = 'Menlo'
        run.font.size = Pt(size)
        run.font.color.rgb = CODE_INK
    return box


# ---------- EDC chrome ---------------------------------------------------
def add_chrome(slide, page_num):
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, Inches(0.18))
    top.fill.solid()
    top.fill.fore_color.rgb = EDC_NAVY
    top.line.fill.background()
    bot = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 0, Inches(7.32), SLIDE_W, Inches(0.18))
    bot.fill.solid()
    bot.fill.fore_color.rgb = EDC_NAVY
    bot.line.fill.background()
    add_text(slide, Inches(0.4), Inches(7.0), Inches(2), Inches(0.3),
             f'{page_num} |', size=10, color=MUTED)
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH,
                                 Inches(12.55), Inches(6.85),
                                 width=Inches(0.65), height=Inches(0.65))


def add_title(slide, bold_part, light_part=None):
    tx = slide.shapes.add_textbox(Inches(0.6), Inches(0.5),
                                  Inches(12.2), Inches(1.0))
    tf = tx.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.margin_right = 0
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.1
    r1 = p.add_run()
    r1.text = bold_part
    r1.font.name = 'Helvetica'
    r1.font.size = Pt(34)
    r1.font.bold = True
    r1.font.color.rgb = EDC_NAVY
    if light_part:
        r2 = p.add_run()
        r2.text = f' | {light_part}'
        r2.font.name = 'Helvetica'
        r2.font.size = Pt(34)
        r2.font.bold = False
        r2.font.color.rgb = EDC_NAVY
    return tx


def set_notes(slide, notes):
    nt = slide.notes_slide
    tf = nt.notes_text_frame
    tf.text = notes


# ---------- Build deck ---------------------------------------------------
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]

SLIDES = []


def slide(fn):
    SLIDES.append(fn)
    return fn


# --- 1. Cover ---
@slide
def s_cover(s):
    fill_bg(s, EDC_NAVY)
    add_text(s, Inches(0), Inches(1.4), SLIDE_W, Inches(0.5),
             'Session 6 · Tirsdag uge 4', size=18, color=EDC_YELLOW,
             align=PP_ALIGN.CENTER, font='Menlo')
    add_text(s, Inches(0), Inches(2.2), SLIDE_W, Inches(1.4),
             'Git, commits, review og PAT', size=54, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(3.5), SLIDE_W, Inches(0.6),
             'Claude i daglig git · branches, commits, reviews · PAT mod Azure DevOps', size=22,
             color=RGBColor(0xC8, 0xD4, 0xE8), align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(5.4), SLIDE_W, Inches(0.5),
             'Præsenteret af Troels Ford Helbo',
             size=20, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(5.95), SLIDE_W, Inches(0.4),
             'Faciliteret af Matias & Michael',
             size=14, italic=True, color=RGBColor(0xC8, 0xD4, 0xE8),
             align=PP_ALIGN.CENTER)
    bot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, Inches(7.32), SLIDE_W, Inches(0.18))
    bot.fill.solid()
    bot.fill.fore_color.rgb = EDC_YELLOW
    bot.line.fill.background()


# --- 2. Dagsorden ---
@slide
def s_agenda(s):
    fill_bg(s, WHITE)
    add_text(s, Inches(0), Inches(0.7), SLIDE_W, Inches(1.0),
             'Dagsorden', size=46, bold=True, color=EDC_NAVY,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(1.6), SLIDE_W, Inches(0.4),
             'Tirsdag · 10:30 – 11:30',
             size=16, color=MUTED, align=PP_ALIGN.CENTER, font='Menlo')
    table_bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), Inches(2.1),
                                  Inches(12.13), Inches(4.7))
    table_bg.fill.solid()
    table_bg.fill.fore_color.rgb = EDC_LIGHT_BG
    table_bg.line.fill.background()

    rows = [
        ('0:00–0:10', 'Recap fra session 5 — hvordan gik subagent-eksperimenterne?', False),
        ('0:10–0:15', 'Næste gang: session 7 — Skills og model-tuning (Oliver)', False),
        ('Fra 0:15', 'Troels præsenterer: Git-workflow, commits og review', True),
        ('Fra 0:40', 'Hands-on — push din første PR til workshop-lab\'en', True),
        ('Til sidst', 'Saml op + take-home + hjemmeopgave', False),
    ]
    row_h = Inches(0.55)
    gap = Inches(0.08)
    start_y = Inches(2.3)
    time_w = Inches(2.0)
    desc_w = Inches(9.6)
    time_x = Inches(0.85)
    desc_x = Inches(2.95)
    for i, (t, w, highlight) in enumerate(rows):
        y = start_y + (row_h + gap) * i
        time_cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, time_x, y, time_w, row_h)
        time_cell.fill.solid()
        time_cell.fill.fore_color.rgb = EDC_NAVY if highlight else EDC_GREY
        time_cell.line.fill.background()
        tt = time_cell.text_frame
        tt.margin_left = Inches(0.2)
        tt.margin_right = Inches(0.2)
        tt.vertical_anchor = MSO_ANCHOR.MIDDLE
        tp = tt.paragraphs[0]
        tp.alignment = PP_ALIGN.LEFT
        tr = tp.add_run()
        tr.text = t
        tr.font.name = 'Helvetica'
        tr.font.size = Pt(15)
        tr.font.bold = True
        tr.font.color.rgb = WHITE
        desc_cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, desc_x, y, desc_w, row_h)
        desc_cell.fill.solid()
        desc_cell.fill.fore_color.rgb = EDC_NAVY if highlight else EDC_GREY
        desc_cell.line.fill.background()
        dt = desc_cell.text_frame
        dt.margin_left = Inches(0.25)
        dt.margin_right = Inches(0.2)
        dt.vertical_anchor = MSO_ANCHOR.MIDDLE
        dp = dt.paragraphs[0]
        dp.alignment = PP_ALIGN.LEFT
        dr = dp.add_run()
        dr.text = w
        dr.font.name = 'Helvetica'
        dr.font.size = Pt(15)
        dr.font.bold = True
        dr.font.color.rgb = WHITE


# --- 3. Recap session 5 ---
@slide
def s_recap(s):
    fill_bg(s, WHITE)
    add_title(s, 'Recap', 'Hvad tog vi med fra session 5?')


# --- 4. Næste gang ---
@slide
def s_next(s):
    fill_bg(s, WHITE)
    add_title(s, 'Næste gang', 'Session 7 · Torsdag uge 4')
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.6), Inches(2.1),
                             Inches(12.13), Inches(3.6))
    box.fill.solid()
    box.fill.fore_color.rgb = EDC_LIGHT_BG
    box.line.fill.background()
    star = s.shapes.add_shape(MSO_SHAPE.STAR_5_POINT,
                              Inches(1.0), Inches(2.5),
                              Inches(1.1), Inches(1.1))
    star.fill.solid()
    star.fill.fore_color.rgb = EDC_YELLOW
    star.line.fill.background()
    add_text(s, Inches(2.6), Inches(2.4), Inches(10), Inches(0.5),
             'SESSION 7 · TORSDAG D. 28. MAJ',
             size=14, bold=True, color=EDC_YELLOW, font='Menlo')
    add_text(s, Inches(2.6), Inches(2.85), Inches(10), Inches(0.7),
             'Skills og model-tuning', size=26, bold=True, color=EDC_NAVY)
    add_rich_text(s, Inches(2.6), Inches(3.55), Inches(10), Inches(0.6),
                  'Custom slash commands · `.claude/skills/` · `/model` og `/effort` · cost vs. speed vs. intelligence',
                  size=16, color=MUTED, line_spacing=1.4)
    add_rich_text(s, Inches(1.0), Inches(4.55), Inches(11.3), Inches(1.0),
                  '**Oliver Pasha Rasoli** præsenterer. Pre-work-link sendes dagen før.',
                  size=15, color=INK, line_spacing=1.5)
    add_text(s, Inches(0.6), Inches(6.0), Inches(12.13), Inches(0.5),
             'Vi mangler stadig præsentanter til session 8 (MCP og hooks) og session 9 (Best practices).',
             size=15, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# --- 5. Section divider: Troels tager scenen ---
@slide
def s_divider_troels(s):
    fill_bg(s, EDC_NAVY)
    add_text(s, Inches(0), Inches(2.7), SLIDE_W, Inches(0.6),
             'Punkt 3 · Dagens emne', size=18, color=EDC_YELLOW,
             align=PP_ALIGN.CENTER, font='Menlo')
    add_text(s, Inches(0), Inches(3.4), SLIDE_W, Inches(1.2),
             'Git + PAT', size=46, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(4.6), SLIDE_W, Inches(0.6),
             'Troels Ford Helbo', size=22,
             color=RGBColor(0xC8, 0xD4, 0xE8), align=PP_ALIGN.CENTER)
    bot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, Inches(7.32), SLIDE_W, Inches(0.18))
    bot.fill.solid()
    bot.fill.fore_color.rgb = EDC_YELLOW
    bot.line.fill.background()


# --- 6. Dagens to ben (HIDDEN backup) ---
@slide
def s_two_pillars(s):
    fill_bg(s, WHITE)
    add_title(s, 'Dagens to ben', 'PAT mod DevOps + Claude i git')
    cards = [
        ('Ben 1 · PAT mod DevOps',
         'Claude læser work items\nfra Azure DevOps via PAT/`az`',
         'Når et boards-ticket beskriver opgaven, kan Claude trække konteksten ind uden copy-paste',
         EDC_NAVY),
        ('Ben 2 · Claude i git',
         'Branches, commits, reviews\n— hele dit daglige flow',
         '`/commit` analyserer diff og rammer jeres stil. `/review` finder rutine-fejl før kollegaen kigger',
         GREEN_INK),
    ]
    col_w = Inches(5.85)
    row_h = Inches(3.6)
    gap_x = Inches(0.3)
    start_x = Inches(0.7)
    start_y = Inches(2.1)
    for i, (title, what, sub, color) in enumerate(cards):
        x = start_x + (col_w + gap_x) * i
        card = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, start_y, col_w, row_h)
        card.adjustments[0] = 0.05
        card.fill.solid()
        card.fill.fore_color.rgb = EDC_LIGHT_BG
        card.line.fill.background()
        stripe = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                    x, start_y, Inches(0.15), row_h)
        stripe.fill.solid()
        stripe.fill.fore_color.rgb = color
        stripe.line.fill.background()
        add_text(s, x + Inches(0.4), start_y + Inches(0.3),
                 col_w - Inches(0.5), Inches(0.6),
                 title, size=24, bold=True, color=color)
        add_text(s, x + Inches(0.4), start_y + Inches(1.1),
                 col_w - Inches(0.5), Inches(1.2),
                 what, size=15, italic=True, color=MUTED, font='Menlo',
                 line_spacing=1.4)
        add_text(s, x + Inches(0.4), start_y + Inches(2.5),
                 col_w - Inches(0.5), Inches(1.0),
                 sub, size=14, color=EDC_NAVY, line_spacing=1.4)
    note = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                              Inches(0.6), Inches(6.0),
                              Inches(12.13), Inches(0.7))
    note.fill.solid()
    note.fill.fore_color.rgb = EDC_NAVY
    note.line.fill.background()
    add_text(s, Inches(0.6), Inches(6.0), Inches(12.13), Inches(0.7),
             'Live demoer til begge ben — vi prøver det selv før vi taler om det.',
             size=14, italic=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# --- 7. Ben 1 · PAT-opsætning (HIDDEN backup) ---
@slide
def s_pat_setup(s):
    fill_bg(s, WHITE)
    add_title(s, 'PAT-opsætning', 'Sådan får Claude adgang til Azure DevOps')
    add_bullets(s, Inches(0.8), Inches(2.1), Inches(11.7), Inches(2.8), [
        'Azure DevOps → **User Settings → Personal Access Tokens** → New Token',
        'Scope: **`Code (read & write)`** + **`Work Items (read)`** så Claude kan trække opgaver',
        'Expiration: **max 90 dage** — sæt en kalender-reminder',
        'Gem i **Git Credential Manager** eller via `az devops login`',
    ], size=17, line_spacing=1.6)
    add_code_block(s, Inches(0.8), Inches(4.9), Inches(11.7), Inches(1.9),
                   '# Variant A — kun git\n'
                   '$ git config --global credential.helper "manager-core"\n'
                   '# Næste git push: indsæt PAT som password → glemt\n\n'
                   '# Variant B — hele Azure-stakken samlet\n'
                   '$ az login && az extension add --name azure-devops && az devops login',
                   size=12)


# --- 8. LIVE DEMO 1: Claude læser DevOps work item (HIDDEN backup) ---
@slide
def s_demo_devops(s):
    fill_bg(s, WHITE)
    add_title(s, 'LIVE DEMO · Ben 1', 'Claude læser en DevOps-opgave')
    add_code_block(s, Inches(0.7), Inches(2.0), Inches(7.8), Inches(4.4),
                   '# Hent et work item via az CLI\n'
                   '$ az boards work-item show --id 12345 \\\n'
                   '    --org https://dev.azure.com/edc \\\n'
                   '    --output json\n\n'
                   '# Eller pump det direkte ind i Claude:\n'
                   '> az boards work-item show --id 12345 \\\n'
                   '    | claude --read\n\n'
                   '# Eller bed Claude selv:\n'
                   '> Læs work item 12345 fra Azure DevOps\n'
                   '  og foreslå hvilke filer der skal ændres\n'
                   '  for at løse acceptkriterierne.',
                   size=13)
    panel = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(8.8), Inches(2.0),
                               Inches(3.95), Inches(4.4))
    panel.adjustments[0] = 0.05
    panel.fill.solid()
    panel.fill.fore_color.rgb = EDC_LIGHT_BG
    panel.line.fill.background()
    add_text(s, Inches(9.0), Inches(2.2), Inches(3.65), Inches(0.4),
             'POINTEN', size=12, bold=True, color=EDC_NAVY, font='Menlo')
    add_text(s, Inches(9.0), Inches(2.65), Inches(3.65), Inches(0.9),
             'Opgaven kommer ind af sig selv',
             size=18, bold=True, color=EDC_NAVY, line_spacing=1.2)
    add_rich_text(s, Inches(9.0), Inches(3.9), Inches(3.65), Inches(2.4),
                  'Ingen copy-paste fra browseren. Claude læser **title, beskrivelse, acceptkriterier, kommentarer** — og ved fra første prompt hvad opgaven handler om.',
                  size=13, color=EDC_NAVY, line_spacing=1.5)
    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
             'Forudsætning: PAT med Work Items (read) + `az` logget ind.',
             size=13, italic=True, color=MUTED, align=PP_ALIGN.CENTER, font='Menlo')


# --- 9. Ben 2 · Claude i daglig git (HIDDEN backup) ---
@slide
def s_git_overview(s):
    fill_bg(s, WHITE)
    add_title(s, 'Claude i daglig git', 'Branches · commits · reviews')
    steps = [
        ('1. Branch',
         '`git checkout -b feature/<id>`\n— eller bed Claude foreslå et navn ud fra DevOps-opgaven',
         EDC_NAVY),
        ('2. Commit',
         '`/commit` analyserer `git diff --cached` og foreslår en besked\ni jeres konvention. Justér én gang — den husker stilen.',
         EDC_BLUE_BODY),
        ('3. Review',
         '`/review` kigger PR\'en igennem FØR kollegaen. Fanger\nnull-checks, glemte tests, scope creep, konvention-brud.',
         GREEN_INK),
    ]
    row_h = Inches(1.2)
    gap = Inches(0.15)
    start_y = Inches(2.0)
    name_w = Inches(2.6)
    body_w = Inches(9.4)
    name_x = Inches(0.7)
    body_x = Inches(3.4)
    for i, (name, body, color) in enumerate(steps):
        y = start_y + (row_h + gap) * i
        name_cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, name_x, y, name_w, row_h)
        name_cell.fill.solid()
        name_cell.fill.fore_color.rgb = color
        name_cell.line.fill.background()
        add_text(s, name_x + Inches(0.25), y, name_w - Inches(0.4), row_h,
                 name, size=22, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        body_cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, body_x, y, body_w, row_h)
        body_cell.fill.solid()
        body_cell.fill.fore_color.rgb = EDC_LIGHT_BG
        body_cell.line.fill.background()
        add_rich_text(s, body_x + Inches(0.3), y, body_w - Inches(0.5), row_h,
                      body, size=14, color=EDC_NAVY,
                      anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)
    add_text(s, Inches(0.7), Inches(6.0), Inches(12), Inches(0.6),
             'Golden rule: Claude-review FORBEREDER kollega-review — det erstatter det ikke.',
             size=15, italic=True, color=EDC_NAVY, align=PP_ALIGN.CENTER, bold=True)


# --- 10. Hvor stoler I på Claude-review (HIDDEN backup) ---
@slide
def s_good_vs_bad(s):
    fill_bg(s, WHITE)
    add_title(s, 'Hvor stoler I på Claude-review?', 'JA og NEJ')
    col_w = Inches(5.85)
    left_x = Inches(0.7)
    right_x = Inches(6.85)
    top_y = Inches(2.0)
    box_h = Inches(4.4)
    good = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left_x, top_y, col_w, box_h)
    good.fill.solid()
    good.fill.fore_color.rgb = GREEN_BG
    good.line.fill.background()
    add_text(s, left_x + Inches(0.4), top_y + Inches(0.3),
             col_w - Inches(0.8), Inches(0.4),
             'JA — god til at fange', size=14, bold=True, color=GREEN_INK, font='Menlo')
    add_bullets(s, left_x + Inches(0.4), top_y + Inches(0.95),
                col_w - Inches(0.8), Inches(3.3), [
        '**Typing**-fejl og null-checks',
        'Konvention-brud (linter ville fange det også)',
        '**Glemte tests** — happy path uden edge cases',
        '**Død kode-stier** der aldrig kaldes',
        'Sløsede imports og ubrugte variabler',
    ], size=14, line_spacing=1.5, bullet_color=GREEN_INK)
    bad = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, right_x, top_y, col_w, box_h)
    bad.fill.solid()
    bad.fill.fore_color.rgb = RED_BG
    bad.line.fill.background()
    add_text(s, right_x + Inches(0.4), top_y + Inches(0.3),
             col_w - Inches(0.8), Inches(0.4),
             'NEJ — kollega skal stadig kigge', size=14, bold=True, color=RED_INK, font='Menlo')
    add_bullets(s, right_x + Inches(0.4), top_y + Inches(0.95),
                col_w - Inches(0.8), Inches(3.3), [
        '**Domæne-logik** — den kender ikke jeres forretning',
        '**Performance** på skala — finder ikke N+1 i prod-data',
        '**Infrastruktur-sikkerhed** — secrets, IAM, network',
        '**Forretningsregler** — "må vi det her ift. compliance?"',
        'Arkitektur-valg der spænder **flere services**',
    ], size=14, line_spacing=1.5, bullet_color=RED_INK)
    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
             'Lad linters tage stil/formatting. Lad Claude tage rutine-fejl. Lad kollegaen tage domænet.',
             size=13, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# --- 11. LIVE DEMO 2: Claude i hele git-flowet (HIDDEN backup) ---
@slide
def s_demo_git(s):
    fill_bg(s, WHITE)
    add_title(s, 'LIVE DEMO · Ben 2', 'Branch → /commit → push → /review')
    add_code_block(s, Inches(0.7), Inches(2.0), Inches(11.93), Inches(4.4),
                   '# 1. Branch (gerne med opgave-ID fra DevOps-work-item)\n'
                   '$ git checkout -b feature/12345-tom-email-guard\n\n'
                   '# 2. Lav ændringen + tests, stage dem\n'
                   '$ git add -p\n\n'
                   '# 3. Claude finder konventionen selv\n'
                   '> /commit\n'
                   '  → "fix(user): guard against empty email in GetUser"\n\n'
                   '# 4. Push og åbn PR i Azure DevOps\n'
                   '$ git push -u origin feature/12345-tom-email-guard\n\n'
                   '# 5. Lad Claude review ÆGTE diff FØR kollegaen\n'
                   '> /review\n'
                   '  → null-checks, glemte tests, scope creep — punktvist',
                   size=13)
    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
             'Hvis Claude justerer commit-stilen forkert: bed om jeres præcise stil én gang — den husker det.',
             size=13, italic=True, color=MUTED, align=PP_ALIGN.CENTER, font='Menlo')


# --- 12. PAT-sikkerhed (HIDDEN backup) ---
@slide
def s_pat_security(s):
    fill_bg(s, WHITE)
    add_title(s, 'PAT-sikkerhed', 'Det her er ikke et password — det er værre')
    cards = [
        ('Aldrig i prompts',
         'PAT i en prompt = PAT i Anthropics logs.\nBrug git credential manager i stedet.',
         RED_INK),
        ('Aldrig i repos',
         'IKKE i `.env` der ligger i git.\nIKKE i Slack/Teams/mails.\nBrug Azure Key Vault eller credential manager.',
         RED_INK),
        ('Min. scope, max. 90 dage',
         '`Code (read & write)` + `Work Items (read)` rækker for de fleste.\nKortere expiration = mindre blast radius.',
         RED_INK),
        ('Revoke straks ved tab',
         'Hvis PAT havner et forkert sted: User Settings → Personal Access Tokens → Revoke.\nLav en ny.',
         RED_INK),
    ]
    row_h = Inches(1.0)
    gap = Inches(0.15)
    start_y = Inches(2.1)
    name_w = Inches(3.3)
    body_w = Inches(8.7)
    name_x = Inches(0.7)
    body_x = Inches(4.1)
    for i, (name, body, color) in enumerate(cards):
        y = start_y + (row_h + gap) * i
        name_cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, name_x, y, name_w, row_h)
        name_cell.fill.solid()
        name_cell.fill.fore_color.rgb = color
        name_cell.line.fill.background()
        add_text(s, name_x + Inches(0.25), y, name_w - Inches(0.4), row_h,
                 name, size=18, bold=True, color=WHITE,
                 anchor=MSO_ANCHOR.MIDDLE)
        body_cell = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, body_x, y, body_w, row_h)
        body_cell.fill.solid()
        body_cell.fill.fore_color.rgb = EDC_LIGHT_BG
        body_cell.line.fill.background()
        add_rich_text(s, body_x + Inches(0.3), y, body_w - Inches(0.5), row_h,
                      body, size=13, color=EDC_NAVY,
                      anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.4)
    add_text(s, Inches(0.7), Inches(6.55), Inches(12), Inches(0.4),
             'PAT giver fuld kode-adgang. Behandl det som et password til hele jeres org\'s repos.',
             size=13, italic=True, color=MUTED, align=PP_ALIGN.CENTER)


# --- 13. GitHub fork-flow til workshop-lab'en (HIDDEN backup) ---
@slide
def s_github_flow(s):
    fill_bg(s, WHITE)
    add_title(s, 'Workshop-lab\'en', 'Hands-on: jeres første PR — fork-flow')
    add_code_block(s, Inches(0.7), Inches(2.0), Inches(11.93), Inches(4.4),
                   '# 1. Fork EDCxClaudeWorkshop på github.com (knap øverst til højre)\n\n'
                   '# 2. Clone din fork — ikke Matias\' repo\n'
                   '$ git clone git@github.com:<dit-handle>/EDCxClaudeWorkshop.git\n'
                   '$ cd EDCxClaudeWorkshop\n'
                   '$ git remote add upstream git@github.com:MatiasGramkow/EDCxClaudeWorkshop.git\n\n'
                   '# 3. Lav en branch til din feature\n'
                   '$ git checkout -b lab/<din-slug>\n\n'
                   '# 4. Kopier eksempel-folderen, rediger din side\n'
                   '$ cp -r demo/app/lab/eksempel demo/app/lab/<din-slug>\n\n'
                   '# 5. Commit + push + åbn PR\n'
                   '$ git add demo/app/lab/<din-slug>/ demo/participants.json\n'
                   '> /commit\n'
                   '$ git push -u origin lab/<din-slug>\n'
                   '# → Vercel poster preview-URL som comment på PR\'en',
                   size=12)


# --- 14. Hands-on divider ---
@slide
def s_divider_handson(s):
    fill_bg(s, EDC_NAVY)
    add_text(s, Inches(0), Inches(2.9), SLIDE_W, Inches(0.6),
             'Punkt 4 · Hands-on', size=22, color=EDC_YELLOW,
             align=PP_ALIGN.CENTER, font='Menlo')
    add_text(s, Inches(0), Inches(3.6), SLIDE_W, Inches(1.2),
             'Push din første PR', size=44, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(4.8), SLIDE_W, Inches(0.6),
             '15 minutter · workshop-lab\'en bliver jeres playground',
             size=20, color=RGBColor(0xC8, 0xD4, 0xE8),
             align=PP_ALIGN.CENTER)
    bot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, Inches(7.32), SLIDE_W, Inches(0.18))
    bot.fill.solid()
    bot.fill.fore_color.rgb = EDC_YELLOW
    bot.line.fill.background()


# --- 15. Hands-on opskrift ---
@slide
def s_handson(s):
    fill_bg(s, WHITE)
    add_title(s, 'Hands-on', 'Push din lab-side live på 15 min')


# --- 16. Cheat-sheet (HIDDEN — Troels har sin egen) ---
@slide
def s_handson_cheat(s):
    fill_bg(s, WHITE)
    add_title(s, 'Cheat-sheet', 'Fork → branch → commit → push → PR')
    a_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(0.7), Inches(2.0),
                               Inches(5.95), Inches(4.7))
    a_box.adjustments[0] = 0.05
    a_box.fill.solid()
    a_box.fill.fore_color.rgb = EDC_LIGHT_BG
    a_box.line.fill.background()
    add_text(s, Inches(0.95), Inches(2.2), Inches(5.5), Inches(0.4),
             'A-SPOR · WORKSHOP LAB', size=12, bold=True, color=EDC_NAVY, font='Menlo')
    add_text(s, Inches(0.95), Inches(2.65), Inches(5.5), Inches(0.6),
             'Push din egen lab-side', size=22, bold=True, color=EDC_NAVY)
    add_bullets(s, Inches(0.95), Inches(3.4), Inches(5.5), Inches(3.2), [
        '**Fork** EDCxClaudeWorkshop på GitHub',
        'Clone din fork lokalt',
        '`cp -r demo/app/lab/eksempel demo/app/lab/<din-slug>`',
        'Brug Claude `/commit`, så `git push`',
        'Åbn PR → få preview-URL på sekunder',
    ], size=13, line_spacing=1.5)

    b_box = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Inches(6.85), Inches(2.0),
                               Inches(5.95), Inches(4.7))
    b_box.adjustments[0] = 0.05
    b_box.fill.solid()
    b_box.fill.fore_color.rgb = CODE_BG
    b_box.line.fill.background()
    add_text(s, Inches(7.1), Inches(2.2), Inches(5.5), Inches(0.4),
             'B-SPOR · AZURE DEVOPS PAT', size=12, bold=True, color=EDC_YELLOW, font='Menlo')
    add_text(s, Inches(7.1), Inches(2.65), Inches(5.5), Inches(0.6),
             'Sæt PAT op til EDC-repos', size=22, bold=True, color=WHITE)
    add_code_block(s, Inches(7.1), Inches(3.4), Inches(5.5), Inches(3.1),
                   '# Opret PAT i Azure DevOps\n'
                   '# (User Settings → PAT → New)\n'
                   '# Scope: Code (read & write)\n\n'
                   '$ git config --global \\\n'
                   '  credential.helper "manager-core"\n\n'
                   '# Næste git pull/push:\n'
                   '# indsæt PAT som password.\n'
                   '# Derefter glemmer du det.',
                   size=11)


# --- 17. Take-home + hjemmeopgave ---
@slide
def s_takehome(s):
    fill_bg(s, WHITE)
    add_title(s, 'Take-home', 'Hjemmeopgave til session 7 — torsdag uge 4')
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             Inches(0.6), Inches(2.0),
                             Inches(12.13), Inches(3.5))
    box.fill.solid()
    box.fill.fore_color.rgb = EDC_LIGHT_BG
    box.line.fill.background()
    add_bullets(s, Inches(1.1), Inches(2.3), Inches(11.2), Inches(3.0), [
        'Brug **`/commit`** eller **`/review`** mindst én gang på ægte arbejde',
        '**Sæt PAT op** til Azure DevOps hvis I ikke har det — så Claude kan push\'e mod EDC-repos',
        'Find **én ting** Claude fangede i et review som I ikke ville have fanget selv',
    ], size=18, line_spacing=1.7)

    bonus = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                               Inches(0.6), Inches(5.7),
                               Inches(12.13), Inches(1.0))
    bonus.fill.solid()
    bonus.fill.fore_color.rgb = EDC_NAVY
    bonus.line.fill.background()
    add_text(s, Inches(0.6), Inches(5.7), Inches(12.13), Inches(1.0),
             'Bonus: tilføj noget til din /lab/<din-slug>-side — vis den frem i session 7',
             size=15, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             font='Menlo')


# --- 18. Tak til Troels ---
@slide
def s_thanks(s):
    fill_bg(s, EDC_NAVY)
    add_text(s, Inches(0), Inches(2.5), SLIDE_W, Inches(1.2),
             'Tak til Troels', size=54, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(3.7), SLIDE_W, Inches(0.6),
             'for at tage session 6', size=24,
             color=RGBColor(0xC8, 0xD4, 0xE8), align=PP_ALIGN.CENTER)
    add_text(s, Inches(0), Inches(5.0), SLIDE_W, Inches(0.5),
             'Næste gang: Skills og model-tuning', size=20,
             color=EDC_YELLOW, align=PP_ALIGN.CENTER, font='Menlo')
    add_text(s, Inches(0), Inches(5.55), SLIDE_W, Inches(0.4),
             'Torsdag d. 28. maj · Oliver Pasha Rasoli', size=14, italic=True,
             color=RGBColor(0xC8, 0xD4, 0xE8), align=PP_ALIGN.CENTER)
    bot = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                             0, Inches(7.32), SLIDE_W, Inches(0.18))
    bot.fill.solid()
    bot.fill.fore_color.rgb = EDC_YELLOW
    bot.line.fill.background()


# ---------- Speaker notes ------------------------------------------------
NOTES = [
    # 1. Cover
    "─── LÆS ───\n"
    "0:00 — alle finder pladser. Matias eller Michael indleder kort, så giver vi ordet til Troels senere.\n\n"
    "─── SIG (Matias/Michael) ───\n"
    "\"Velkommen til session 6. I dag er det Troels der præsenterer — to ben: PAT mod Azure DevOps så Claude kan læse work items, og hvordan Claude bruger git med branches, commits og reviews. Live demoer til begge.\"\n\n"
    "\"Efter Troels prøver I selv at pushe en PR til workshop-lab'en.\"",

    # 2. Dagsorden
    "─── LÆS ───\n"
    "0:00–0:02 — vis flowet. Peg på de fremhævede (blå) felter — det er hvor det meste af tiden går.\n\n"
    "─── SIG ───\n"
    "\"Samme struktur: recap, dagens emne, hands-on, take-home.\"\n\n"
    "\"Hands-on i dag er konkret: I forker, cloner, laver en branch, kører `/commit`, pusher, og åbner PR. Når det er gjort har I jeres egen URL på workshop-sitet.\"",

    # 3. Recap session 5
    "─── LÆS ───\n"
    "0:02–0:10 — kort runde. Hjemmeopgaven var: kør Explore på et EDC-repo, byg én custom subagent, spawn 2+ parallelle. Tag 2-3 deltagere.\n\n"
    "─── SIG ───\n"
    "\"Hjemmeopgaven: Explore + én custom agent + parallelle subagents. Hurtig runde — hvad byggede I?\"\n\n"
    "─── GØR ───\n"
    "Hold tiden. Max 8 min på recap så Troels får sine 25 min.",

    # 4. Næste gang
    "─── LÆS ───\n"
    "0:10–0:15 — kort annonce. Oliver er allerede på næste session, så det her er bare en heads-up, ikke en rekruttering.\n\n"
    "─── SIG ───\n"
    "\"Inden Troels får ordet — heads-up om næste gang: Oliver Pasha Rasoli tager session 7 om Skills og model-tuning, torsdag den 28. maj. I får pre-work dagen før.\"\n\n"
    "\"Vi mangler stadig præsentanter til session 8 og 9 — kig forbi mig eller Michael bagefter hvis I vil tage en.\"",

    # 5. Divider — Troels tager over
    "─── LÆS ───\n"
    "0:15 — overdrag scenen. Sig Troels' navn, klap kort, sæt jer ned. Troels har ordet de næste 25 minutter — OG hands-on bagefter.\n\n"
    "─── SIG (Matias/Michael) ───\n"
    "\"Troels — værsgo, scenen er din.\"\n\n"
    "─── GØR ───\n"
    "Sæt jer på sidelinjen. Lad Troels køre. Hop kun ind hvis han direkte beder om det.",

    # 6. Dagens to ben (HIDDEN backup)
    "[SKJULT BACKUP — Troels bygger sine egne slides. Disse ligger i decket hvis han skal bruge dem.]\n\n"
    "─── HVIS VIST (Troels) ───\n"
    "Sæt rammen i 30 sekunder: dagens to ben. PAT mod DevOps (læser opgaver), og Claude i git (branches, commits, reviews). Live demoer til begge — vi prøver det selv før vi taler om det.",

    # 7. PAT-opsætning (HIDDEN backup)
    "[SKJULT BACKUP]\n\n"
    "─── HVIS VIST (Troels) ───\n"
    "Det her er hvor de fleste sidder fast: 'Claude kan ikke pushe til mine repos eller læse opgaver.' Svaret er PAT.\n\n"
    "Gå roligt igennem: New Token → Scope `Code (read & write)` + `Work Items (read)` → Expiration max 90 dage. To opsætningsveje: Git Credential Manager til kun-git, eller `az devops login` hvis I bruger hele Azure-stakken.",

    # 8. LIVE DEMO · Ben 1 (HIDDEN backup)
    "[SKJULT BACKUP — LIVE DEMO 1]\n\n"
    "─── HVIS VIST (Troels) ───\n"
    "Live demo: vis `az boards work-item show --id <ID>` på et rigtigt EDC work item. Vis JSON-output kort, og pump det så ind til Claude med 'Læs work item X og foreslå filer at ændre'.\n\n"
    "Pointe: ingen copy-paste fra browser. Claude læser title, beskrivelse, acceptkriterier og kommentarer direkte. Den ved fra første prompt hvad opgaven handler om.\n\n"
    "Hvis demoen fejler (PAT ikke sat op, work item ikke fundet): fall back til JSON-eksemplet og forklar konceptet.",

    # 9. Ben 2 · Claude i daglig git (HIDDEN backup)
    "[SKJULT BACKUP]\n\n"
    "─── HVIS VIST (Troels) ───\n"
    "Tre trin: branch → commit → review. Pointer:\n"
    "- Branch-navn kan komme fra DevOps-opgaven (fx 'feature/12345-tom-email-guard')\n"
    "- `/commit` finder konvention fra seneste commits — første gang du retter stilen, husker den\n"
    "- `/review` er FØR kollegaen — ikke i stedet for",

    # 10. JA/NEJ — Claude-review (HIDDEN backup)
    "[SKJULT BACKUP]\n\n"
    "─── HVIS VIST (Troels) ───\n"
    "Den vigtigste pointe i hele sessionen: Claude er god til det rutineprægede. Den er IKKE god til domænet. Lad jeres kollega tage domænet.\n\n"
    "Sig højt: 'Claude-review forbereder kollega-review — det erstatter det ikke.'",

    # 11. LIVE DEMO · Ben 2 (HIDDEN backup)
    "[SKJULT BACKUP — LIVE DEMO 2]\n\n"
    "─── HVIS VIST (Troels) ───\n"
    "Live demo: hele git-flowet i ét take.\n"
    "1. Lav branch (gerne med DevOps work item-ID i navnet — bygger bro tilbage til Ben 1)\n"
    "2. Stage en ægte ændring (du har en klar i forvejen)\n"
    "3. `/commit` — vis at den fanger stilen. Hvis ikke: bed om jeres præcise stil og lad den rette sig\n"
    "4. Push branchen\n"
    "5. `/review` — gennemgå punktvist højt: 'ja, ret' / 'nej, ikke relevant' / 'udenfor scope'\n\n"
    "Hvis demoen tager for lang tid: spring trin 4 over og kør `/review` på en branch du allerede har klar.",

    # 12. PAT-sikkerhed (HIDDEN backup)
    "[SKJULT BACKUP]\n\n"
    "─── HVIS VIST (Troels) ───\n"
    "PAT er ikke et almindeligt password. Det er fuld kode-adgang til alt I har rettigheder til.\n\n"
    "Hvis I tror den er lækket: revoke STRAKS. Lav en ny. Bedre at have 5 minutters PAT-pause end at miste søvn over hvor den ligger.",

    # 13. GitHub fork-flow (HIDDEN backup)
    "[SKJULT BACKUP]\n\n"
    "─── HVIS VIST (Troels) ───\n"
    "Vis det her hvis nogen sidder fast i workshop-lab-flowet. Tre vigtige punkter:\n"
    "1. Fork-knappen er øverst til højre på GitHub\n"
    "2. Clone DIN fork, ikke Matias' repo\n"
    "3. `git remote add upstream` så I kan hente fra Matias' senere\n\n"
    "Resten er almindelig branch + commit + push + PR.",

    # 14. Hands-on divider
    "─── LÆS (Troels) ───\n"
    "0:40 — overgang til hands-on. Troels driver det selv. Energien skal op — det her er hvor de bygger noget de kan vise.\n\n"
    "─── SIG (Troels) ───\n"
    "\"Nu er det jer. 15 minutter — push jeres første PR til workshop-lab'en.\"\n\n"
    "─── GØR ───\n"
    "Klap én gang, gå videre til næste slide hvor opskriften er.",

    # 15. Hands-on opskrift
    "─── LÆS (Troels) ───\n"
    "0:40–0:55 — opskriften står her. Matias og Michael går rundt og hjælper. Hvis nogen sidder fast, bed dem råbe.\n\n"
    "─── SIG (Troels) ───\n"
    "\"Find en partner. Vælg A-spor hvis I vil pushe til workshop-lab'en, B-spor hvis I først vil have PAT sat op til Azure DevOps. I behøver ikke samme spor.\"\n\n"
    "\"Hele flowet ligger i demo/LAB.md i repoet — Claude læser den automatisk hvis I beder den hjælpe.\"\n\n"
    "─── GØR ───\n"
    "Lad næste slide stå mens de arbejder — cheat-sheet er nemmere at glo på end at huske udenad.",

    # 16. Cheat-sheet (HIDDEN — Troels har sin egen)
    "[SKJULT BACKUP — Troels bringer sin egen cheat-sheet hvis han vil have en.]\n\n"
    "─── HVIS VIST (Troels) ───\n"
    "Skift hertil straks de begynder. Lad det stå hele hands-on. De skal kunne kigge op og se kommandoerne.\n\n"
    "─── GØR ───\n"
    "Gå rundt med Matias og Michael. Læg mærke til den første der får preview-URL — bed dem dele den når I samler op.",

    # 17. Take-home
    "─── LÆS (Troels/Matias) ───\n"
    "0:55–1:00 — saml op. Spørg højt hvem der fik en preview-URL. Vis 1-2 af dem på storskærm hvis muligt. Så hjemmeopgaven.\n\n"
    "─── SIG ───\n"
    "\"Hvem fik en preview-URL? Lad os se én eller to.\"\n\n"
    "[Vis dem på storskærm hvis muligt.]\n\n"
    "\"Hjemmeopgaven: brug `/commit` eller `/review` på ægte arbejde mindst én gang, sæt PAT op hvis I ikke har det, og find én ting Claude fangede i et review.\"\n\n"
    "\"Bonus: byg videre på jeres /lab-side — vis den frem næste gang.\"",

    # 18. Tak
    "─── LÆS (Matias/Michael) ───\n"
    "1:00 — afrund. Tak Troels højt foran rummet.\n\n"
    "─── SIG ───\n"
    "\"Stort tak til Troels for at tage session 6. Næste gang er session 7 — Skills og model-tuning — torsdag den 28. maj med Oliver.\"\n\n"
    "\"Tak for i dag.\"",
]


# ---------- Build deck ---------------------------------------------------
total = len(SLIDES)
# Cover (0), Troels-divider (4), Hands-on-divider (13) og Tak (17)
# har egen full-bleed baggrund — ingen chrome.
DIVIDER_INDEXES = {0, 4, 13, 17}
# Slides 6-13 (Troels' indhold) og slide 16 (cheat-sheet) er skjult —
# Troels præsenterer sit eget materiale. Skjulte slides ligger som backup.
HIDDEN_INDEXES = set(range(5, 13)) | {15}

for idx, builder in enumerate(SLIDES):
    s = prs.slides.add_slide(blank)
    builder(s)
    page_num = idx + 1
    if idx not in DIVIDER_INDEXES:
        add_chrome(s, page_num)
    if idx in HIDDEN_INDEXES:
        s.element.set('show', '0')
    if idx < len(NOTES):
        set_notes(s, NOTES[idx])

out = os.path.join(HERE, 'session6.pptx')
prs.save(out)
print(f'Wrote {out}  ({total} slides)')
